#!/usr/bin/env python3
"""Build ERD + code-first migrations lecture deck (Google Slides–compatible PPTX)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT_DIR = Path(__file__).resolve().parent
OUT_FILE = OUT_DIR / "erd-migrations-lecture.pptx"

NAVY = RGBColor(0x0F, 0x17, 0x2A)
TEAL = RGBColor(0x0D, 0x94, 0x88)
TEAL_LIGHT = RGBColor(0x14, 0xB8, 0xA6)
SLATE = RGBColor(0x47, 0x55, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF8, 0xFA, 0xFC)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
BLUE_PANEL = RGBColor(0xE3, 0xF2, 0xFD)

FONT_TITLE = "Poppins"
FONT_BODY = "Lato"
FONT_MONO = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_round_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=18,
    bold=False,
    color=NAVY,
    align=PP_ALIGN.LEFT,
    font=FONT_BODY,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return box, tf


def add_bullets(tf, items, size=16, color=SLATE, spacing=6, level=0):
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = level
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT_BODY
        p.space_after = Pt(spacing)
        p.bullet = True


def add_header_bar(slide, kicker=None):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12), TEAL)
    if kicker:
        add_textbox(
            slide,
            Inches(0.7),
            Inches(0.35),
            Inches(5),
            Inches(0.35),
            kicker.upper(),
            size=11,
            bold=True,
            color=TEAL,
            font=FONT_TITLE,
        )


def add_footer(slide, text="PayNest · Live session 2026-07-16"):
    add_textbox(
        slide, Inches(0.7), Inches(7.05), Inches(8), Inches(0.3), text, size=9, color=SLATE
    )


def slide_title(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, TEAL)
    add_round_rect(slide, Inches(9.2), Inches(1.2), Inches(3.6), Inches(5.1), TEAL_LIGHT)
    add_round_rect(slide, Inches(9.6), Inches(1.6), Inches(2.8), Inches(4.3), TEAL)
    add_textbox(
        slide, Inches(0.9), Inches(1.5), Inches(7.8), Inches(0.5),
        "LIVE SESSION", size=14, bold=True, color=TEAL_LIGHT, font=FONT_TITLE,
    )
    add_textbox(
        slide, Inches(0.9), Inches(2.0), Inches(8), Inches(1.8),
        "ERD + Code-First\nMigrations",
        size=40, bold=True, color=WHITE, font=FONT_TITLE,
    )
    add_textbox(
        slide, Inches(0.9), Inches(4.0), Inches(7.5), Inches(1.0),
        "From Java domain classes → SQL migrations → H2 database (Capstone 4 prep)",
        size=20, color=OFF_WHITE,
    )
    add_textbox(
        slide, Inches(0.9), Inches(6.2), Inches(7), Inches(0.4),
        "First-time programmers · Plain JDBC · No ORM",
        size=13, color=SLATE,
    )
    slide.notes_slide.notes_text_frame.text = (
        "Welcome. Full walkthrough: 01-code-first-migrations-walkthrough.md. "
        "ERD is visualisation only; migrations come from Java classes."
    )


def slide_problem(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Why a database?")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "Memory is not durable", size=34, bold=True, color=NAVY, font=FONT_TITLE,
    )
    _, tf = add_textbox(slide, Inches(0.7), Inches(1.75), Inches(11.5), Inches(2.5), "", size=17)
    add_bullets(
        tf,
        [
            "In-memory stores (HashMap) lose everything on restart or crash.",
            "Finance needs counts that match persisted rows — not developer RAM.",
            "Retries: “Did we charge twice?” — you need a row to point at.",
            "Capstone 4: payment attempts must survive restarts using H2 on disk.",
        ],
        size=17,
    )
    banner = add_round_rect(slide, Inches(0.7), Inches(4.8), Inches(11.9), Inches(1.5), TEAL)
    add_textbox(
        slide, Inches(1.0), Inches(5.05), Inches(11.3), Inches(1.0),
        "Database = long-term storage on disk. Java reloads rows after restart.",
        size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    add_footer(slide)


def slide_session_spine(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Overview")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "Session spine", size=34, bold=True, color=NAVY, font=FONT_TITLE,
    )
    steps = [
        ("Step 1", "Understand domain classes → map fields to columns"),
        ("Step 2", "Create migrations (write CREATE TABLE SQL)"),
        ("Step 3", "Apply migration to H2 via JDBC"),
        ("Step 4", "Preview: Java stores INSERT / SELECT"),
        ("Step 5", "Reset schema (wipe local H2 runbook)"),
    ]
    for i, (label, body) in enumerate(steps):
        top = Inches(1.9 + i * 1.05)
        add_round_rect(slide, Inches(0.7), top, Inches(11.9), Inches(0.9), WHITE)
        add_rect(slide, Inches(0.7), top, Inches(0.08), Inches(0.9), TEAL)
        add_textbox(
            slide, Inches(1.0), top + Inches(0.12), Inches(1.5), Inches(0.5),
            label, size=16, bold=True, color=TEAL, font=FONT_TITLE,
        )
        add_textbox(
            slide, Inches(2.5), top + Inches(0.12), Inches(9.5), Inches(0.6),
            body, size=16, color=SLATE,
        )
    add_textbox(
        slide, Inches(0.7), Inches(6.35), Inches(11.5), Inches(0.5),
        "Source of truth: Java classes in com.paynestsystem.domain — not the ERD diagram.",
        size=14, bold=True, color=AMBER, font=FONT_TITLE,
    )
    add_footer(slide)


def slide_erd_visual_only(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "ERD")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "ERD = visualisation only", size=34, bold=True, color=NAVY, font=FONT_TITLE,
    )
    panel = add_round_rect(slide, Inches(0.7), Inches(1.85), Inches(5.8), Inches(4.8), WHITE)
    _, tf = add_textbox(slide, Inches(1.0), Inches(2.1), Inches(5.3), Inches(4.3), "", size=16)
    add_bullets(
        tf,
        [
            "paynest-erd.drawio helps you see entities & relationships.",
            "Open on projector or while studying — optional.",
            "Do NOT export SQL from draw.io.",
            "If stuck on relationships: peek at diagram, then return to .java files.",
        ],
        size=16,
    )
    panel2 = add_round_rect(slide, Inches(6.8), Inches(1.85), Inches(5.8), Inches(4.8), BLUE_PANEL)
    add_textbox(
        slide, Inches(7.1), Inches(2.1), Inches(5.2), Inches(0.5),
        "Do this", size=18, bold=True, color=NAVY, font=FONT_TITLE,
    )
    _, tf2 = add_textbox(slide, Inches(7.1), Inches(2.65), Inches(5.2), Inches(2.0), "", size=15)
    add_bullets(tf2, ["Write SQL from Java fields", "Use diagram to check understanding"], size=15)
    add_textbox(
        slide, Inches(7.1), Inches(4.2), Inches(5.2), Inches(0.5),
        "Do not do this", size=18, bold=True, color=NAVY, font=FONT_TITLE,
    )
    _, tf3 = add_textbox(slide, Inches(7.1), Inches(4.75), Inches(5.2), Inches(1.5), "", size=15)
    add_bullets(tf3, ["Wait for ERD before Step 2", "Copy column names only from boxes"], size=15)
    add_footer(slide)


def slide_glossary_1(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Glossary")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "Key terms (1 of 2)", size=34, bold=True, color=NAVY, font=FONT_TITLE,
    )
    terms = [
        ("Entity", "A thing you care about → Java class → SQL table"),
        ("Column / field", "One piece of data (e.g. Customer.name → customers.name)"),
        ("Primary key (PK)", "Uniquely identifies one row (like a student number)"),
        ("Foreign key (FK)", "Points to another table's PK (orders.customer_id → customers.id)"),
        ("Migration", "Versioned SQL that changes DB structure (schema-v1.sql)"),
        ("DDL vs DML", "DDL = CREATE TABLE; DML = INSERT / SELECT after tables exist"),
    ]
    for i, (term, defn) in enumerate(terms):
        col = i % 2
        row = i // 2
        left = Inches(0.7 + col * 6.2)
        top = Inches(1.85 + row * 1.75)
        add_round_rect(slide, left, top, Inches(5.9), Inches(1.55), WHITE)
        add_textbox(
            slide, left + Inches(0.25), top + Inches(0.15), Inches(5.4), Inches(0.4),
            term, size=15, bold=True, color=TEAL, font=FONT_TITLE,
        )
        add_textbox(
            slide, left + Inches(0.25), top + Inches(0.55), Inches(5.4), Inches(0.85),
            defn, size=14, color=SLATE,
        )
    add_footer(slide)


def slide_glossary_2(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Glossary")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "Key terms (2 of 2)", size=34, bold=True, color=NAVY, font=FONT_TITLE,
    )
    terms = [
        ("H2", "Small Java DB; file URL: jdbc:h2:file:./data/paynest"),
        ("JDBC", "Java API to connect and run SQL (DriverManager, Connection)"),
        ("Code-first", "Classes exist first → you author DDL to match"),
        ("Idempotency key", "Same key on retry = same outcome; store as UNIQUE"),
        ("ERD", "Picture of the model — not a build step for PayNest"),
        ("No EF / Flyway here", "You write SQL manually; Capstone 4 wants plain JDBC"),
    ]
    for i, (term, defn) in enumerate(terms):
        col = i % 2
        row = i // 2
        left = Inches(0.7 + col * 6.2)
        top = Inches(1.85 + row * 1.75)
        add_round_rect(slide, left, top, Inches(5.9), Inches(1.55), WHITE)
        add_textbox(
            slide, left + Inches(0.25), top + Inches(0.15), Inches(5.4), Inches(0.4),
            term, size=15, bold=True, color=TEAL, font=FONT_TITLE,
        )
        add_textbox(
            slide, left + Inches(0.25), top + Inches(0.55), Inches(5.4), Inches(0.85),
            defn, size=14, color=SLATE,
        )
    add_footer(slide)


def slide_code_first(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12), TEAL_LIGHT)
    add_textbox(
        slide, Inches(0.7), Inches(0.5), Inches(11), Inches(0.8),
        "Code-first (PayNest)", size=34, bold=True, color=WHITE, font=FONT_TITLE,
    )
    flow = [
        "Java domain classes already exist",
        "Decide what must survive restart (Capstone 4: payments)",
        "Step 2: WRITE CREATE TABLE SQL",
        "Step 3: APPLY SQL to H2 with JDBC",
        "Capstone 4: stores that INSERT / SELECT",
    ]
    for i, line in enumerate(flow):
        top = Inches(1.6 + i * 1.0)
        add_textbox(
            slide, Inches(1.2), top, Inches(0.5), Inches(0.5),
            "→", size=22, bold=True, color=TEAL_LIGHT,
        )
        add_textbox(slide, Inches(1.8), top, Inches(10), Inches(0.55), line, size=18, color=OFF_WHITE)
    add_round_rect(slide, Inches(0.7), Inches(6.0), Inches(11.9), Inches(0.75), TEAL)
    add_textbox(
        slide, Inches(1.0), Inches(6.15), Inches(11.3), Inches(0.45),
        "Database-first (tables first, classes later) — not our path.",
        size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    add_footer(slide, "No Spring ddl-auto · No Hibernate in this scaffold")


def slide_step1_intro(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Step 1")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "Understand the domain (classes)", size=34, bold=True, color=NAVY, font=FONT_TITLE,
    )
    add_textbox(
        slide, Inches(0.7), Inches(1.65), Inches(11), Inches(0.5),
        "Goal: know fields → decide tables & columns before writing SQL.",
        size=18, color=SLATE,
    )
    add_round_rect(slide, Inches(0.7), Inches(2.4), Inches(5.5), Inches(3.8), WHITE)
    add_textbox(
        slide, Inches(1.0), Inches(2.65), Inches(5), Inches(0.4),
        "Commerce (Capstone 1)", size=16, bold=True, color=TEAL, font=FONT_TITLE,
    )
    _, tf = add_textbox(slide, Inches(1.0), Inches(3.1), Inches(4.9), Inches(2.8), "", size=15)
    add_bullets(
        tf,
        [
            "Product — id, name, price",
            "Customer — id, name, email",
            "Order — id, customer (object!), items",
            "OrderItem — product, quantity",
            "Usually in-memory today — practise mapping",
        ],
        size=15,
    )
    add_round_rect(slide, Inches(6.5), Inches(2.4), Inches(6.1), Inches(3.8), BLUE_PANEL)
    add_textbox(
        slide, Inches(6.8), Inches(2.65), Inches(5.5), Inches(0.4),
        "Payments (Capstone 4–5)", size=16, bold=True, color=NAVY, font=FONT_TITLE,
    )
    _, tf2 = add_textbox(slide, Inches(6.8), Inches(3.1), Inches(5.5), Inches(2.8), "", size=15)
    add_bullets(
        tf2,
        [
            "Transaction — amount, bank, timestamp",
            "TransactionRecord — id, idempotencyKey, nested Transaction, status…",
            "AiDecisionRecord — audit trail (Capstone 5)",
            "These must become durable H2 tables",
        ],
        size=15,
    )
    add_footer(slide)


def slide_step1_nested(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Step 1")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "Nested objects → flatten or FK", size=34, bold=True, color=NAVY, font=FONT_TITLE,
    )
    add_textbox(
        slide, Inches(0.7), Inches(1.65), Inches(11.5), Inches(0.55),
        "Databases don't store Java objects inside a column the same way.",
        size=18, color=SLATE,
    )
    examples = [
        (
            "Order.customer",
            "Becomes orders.customer_id (FK → customers.id)",
        ),
        (
            "Transaction inside TransactionRecord",
            "Flatten: amount, bank, transaction_timestamp on transaction_records",
        ),
        (
            "Java enum (TransactionStatus)",
            "Store as VARCHAR: PENDING, ROUTED, COMPLETED, FAILED",
        ),
    ]
    for i, (head, sub) in enumerate(examples):
        top = Inches(2.4 + i * 1.45)
        add_round_rect(slide, Inches(0.7), top, Inches(11.9), Inches(1.25), WHITE)
        add_textbox(
            slide, Inches(1.0), top + Inches(0.2), Inches(4.5), Inches(0.45),
            head, size=16, bold=True, color=TEAL, font=FONT_MONO,
        )
        add_textbox(
            slide, Inches(5.5), top + Inches(0.2), Inches(6.8), Inches(0.8),
            sub, size=15, color=SLATE,
        )
    add_footer(slide)


def slide_step1_mapping(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Step 1")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "Worked example: TransactionRecord", size=34, bold=True, color=NAVY, font=FONT_TITLE,
    )
    rows = [
        ("id", "id", "VARCHAR", "PK"),
        ("idempotencyKey", "idempotency_key", "VARCHAR", "UNIQUE"),
        ("transaction.amount", "amount", "DOUBLE", "flattened"),
        ("transaction.bank", "bank", "VARCHAR", "flattened"),
        ("transaction.timestamp", "transaction_timestamp", "TIMESTAMP", "flattened"),
        ("status", "status", "VARCHAR", "enum as text"),
    ]
    headers = ["Java field", "SQL column", "Type", "Notes"]
    table = slide.shapes.add_table(len(rows) + 1, 4, Inches(0.7), Inches(1.75), Inches(11.9), Inches(4.5))
    tbl = table.table
    col_widths = [Inches(3.2), Inches(3.4), Inches(2.0), Inches(3.3)]
    for j, w in enumerate(col_widths):
        tbl.columns[j].width = w
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = val
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].font.name = FONT_MONO if j < 2 else FONT_BODY
    add_textbox(
        slide, Inches(0.7), Inches(6.4), Inches(11.5), Inches(0.45),
        "Build this mapping sheet from .java files — before opening draw.io.",
        size=14, bold=True, color=AMBER, font=FONT_TITLE,
    )
    add_footer(slide)


def slide_step1_checkpoint(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, TEAL)
    add_textbox(
        slide, Inches(0.7), Inches(0.9), Inches(11), Inches(0.8),
        "Checkpoint — Step 1", size=34, bold=True, color=WHITE, font=FONT_TITLE,
    )
    _, tf = add_textbox(slide, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.5), "", size=20, color=WHITE)
    add_bullets(
        tf,
        [
            "Name two Java classes Capstone 4 needs durable.",
            "Why is idempotency_key special?",
            "What happens to Transaction fields in transaction_records?",
        ],
        size=22,
        color=WHITE,
        spacing=18,
    )
    add_textbox(
        slide, Inches(0.7), Inches(6.2), Inches(11), Inches(0.4),
        "Answer without opening draw.io.",
        size=14, color=OFF_WHITE,
    )
    slide.notes_slide.notes_text_frame.text = "Pause for discussion. Expected: TransactionRecord, idempotency; flatten Transaction columns."


def slide_step2_intro(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Step 2")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "Create the migrations (DDL)", size=34, bold=True, color=NAVY, font=FONT_TITLE,
    )
    add_textbox(
        slide, Inches(0.7), Inches(1.65), Inches(11.5), Inches(0.55),
        "Turn Step 1 mapping into CREATE TABLE SQL. Nothing auto-generates this for PayNest.",
        size=18, color=SLATE,
    )
    _, tf = add_textbox(slide, Inches(0.7), Inches(2.4), Inches(5.8), Inches(3.5), "", size=16)
    add_bullets(
        tf,
        [
            "Work in schema-v1.sql (this folder)",
            "Copy into H2Schema.java for Capstone deliverable",
            "Required: transaction_records, ai_decision_records",
            "Optional: commerce tables (commented in schema-v1.sql)",
        ],
        size=16,
    )
    code_box = add_round_rect(slide, Inches(6.8), Inches(2.4), Inches(5.8), Inches(3.8), NAVY)
    add_textbox(
        slide, Inches(7.0), Inches(2.55), Inches(5.4), Inches(3.5),
        "No dotnet ef migrations add\nNo Hibernate ddl-auto\nNo draw.io export\n\nYou author SQL from Java fields.",
        size=15, color=OFF_WHITE, font=FONT_MONO,
    )
    add_footer(slide)


def slide_step2_sql_keywords(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Step 2")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "SQL keywords — why they exist", size=34, bold=True, color=NAVY, font=FONT_TITLE,
    )
    kw = [
        ("CREATE TABLE", "New place to store rows"),
        ("IF NOT EXISTS", "Safe to re-run while learning"),
        ("NOT NULL", "Column must have a value"),
        ("PRIMARY KEY", "One unique id per row"),
        ("UNIQUE", "No duplicates — idempotency safety net"),
        ("FOREIGN KEY", "Must match parent table's PK"),
        ("CREATE INDEX", "Faster lookups (reports, keys)"),
    ]
    for i, (k, v) in enumerate(kw):
        col = i % 2
        row = i // 2
        left = Inches(0.7 + col * 6.2)
        top = Inches(1.85 + row * 1.45)
        add_round_rect(slide, left, top, Inches(5.9), Inches(1.25), WHITE)
        add_textbox(
            slide, left + Inches(0.25), top + Inches(0.15), Inches(5.4), Inches(0.4),
            k, size=14, bold=True, color=TEAL, font=FONT_MONO,
        )
        add_textbox(
            slide, left + Inches(0.25), top + Inches(0.55), Inches(5.4), Inches(0.55),
            v, size=14, color=SLATE,
        )
    add_footer(slide)


def slide_step2_create_table(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Step 2")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "CREATE TABLE example", size=34, bold=True, color=NAVY, font=FONT_TITLE,
    )
    sql = """CREATE TABLE IF NOT EXISTS transaction_records (
  id VARCHAR(64) NOT NULL,
  idempotency_key VARCHAR(128) NOT NULL,
  amount DOUBLE NOT NULL,
  bank VARCHAR(64) NOT NULL,
  transaction_timestamp TIMESTAMP NOT NULL,
  status VARCHAR(32) NOT NULL,
  created_at TIMESTAMP NOT NULL,
  updated_at TIMESTAMP NOT NULL,
  CONSTRAINT pk_transaction_records
    PRIMARY KEY (id),
  CONSTRAINT uq_idempotency
    UNIQUE (idempotency_key)
);"""
    add_round_rect(slide, Inches(0.7), Inches(1.75), Inches(11.9), Inches(4.9), NAVY)
    add_textbox(
        slide, Inches(0.95), Inches(1.95), Inches(11.4), Inches(4.5),
        sql, size=13, color=OFF_WHITE, font=FONT_MONO,
    )
    add_textbox(
        slide, Inches(0.7), Inches(6.75), Inches(11.5), Inches(0.35),
        "Full version with comments: schema-v1.sql in this folder.",
        size=13, color=SLATE,
    )
    add_footer(slide)


def slide_step2_mistakes(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Step 2")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "Common mistakes", size=34, bold=True, color=NAVY, font=FONT_TITLE,
    )
    _, tf = add_textbox(slide, Inches(0.7), Inches(1.85), Inches(11.5), Inches(4.5), "", size=17)
    add_bullets(
        tf,
        [
            "Waiting for the ERD before writing SQL — classes are enough.",
            "Forgetting UNIQUE on idempotency_key — duplicate retries could double-post.",
            "Storing a Java object in one column — flatten or use FK.",
            "Spaces in table names — use snake_case (transaction_records).",
            "Naming columns from the diagram without checking .java files.",
        ],
        size=17,
    )
    add_footer(slide)


def slide_step2_checkpoint(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, TEAL)
    add_textbox(
        slide, Inches(0.7), Inches(0.9), Inches(11), Inches(0.8),
        "Checkpoint — Step 2", size=34, bold=True, color=WHITE, font=FONT_TITLE,
    )
    _, tf = add_textbox(slide, Inches(0.9), Inches(2.0), Inches(11.5), Inches(4.5), "", size=20, color=WHITE)
    add_bullets(
        tf,
        [
            "What is a migration in your own words?",
            "Why UNIQUE (idempotency_key)?",
            "Where will this SQL live in the repo for Capstone 4?",
        ],
        size=22,
        color=WHITE,
        spacing=18,
    )
    slide.notes_slide.notes_text_frame.text = "H2Schema.java or schema artefact. UNIQUE prevents duplicate business attempts."


def slide_step3_apply(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Step 3")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "Apply the migration", size=34, bold=True, color=NAVY, font=FONT_TITLE,
    )
    add_textbox(
        slide, Inches(0.7), Inches(1.65), Inches(11.5), Inches(0.5),
        "jdbc:h2:file:./data/paynest  →  files under ./data/ (gitignored)",
        size=17, bold=True, color=TEAL, font=FONT_MONO,
    )
    flow = [
        "SQL text (schema-v1.sql or H2Schema constant)",
        "Java opens JDBC Connection",
        "Execute CREATE TABLE statements",
        "H2 writes files to ./data/",
        "Tables exist — empty until INSERT",
    ]
    for i, line in enumerate(flow):
        top = Inches(2.35 + i * 0.85)
        add_textbox(slide, Inches(1.0), top, Inches(0.4), Inches(0.4), "↓", size=18, color=TEAL)
        add_textbox(slide, Inches(1.5), top, Inches(10.5), Inches(0.5), line, size=16, color=SLATE)
    add_round_rect(slide, Inches(0.7), Inches(6.15), Inches(11.9), Inches(0.65), AMBER)
    add_textbox(
        slide, Inches(1.0), Inches(6.28), Inches(11.3), Inches(0.4),
        "Run DDL once at startup — not on every save.",
        size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    add_footer(slide)


def slide_step3_jdbc(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Step 3")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "JDBC sketch (learning only)", size=34, bold=True, color=NAVY, font=FONT_TITLE,
    )
    java = """Connection conn = DriverManager.getConnection(
    "jdbc:h2:file:./data/paynest", "sa", "");
try (conn; Statement stmt = conn.createStatement()) {
    stmt.execute(H2Schema.CREATE_TRANSACTION_RECORDS);
    // ... other DDL ...
}"""
    add_round_rect(slide, Inches(0.7), Inches(1.75), Inches(11.9), Inches(3.2), NAVY)
    add_textbox(
        slide, Inches(0.95), Inches(1.95), Inches(11.4), Inches(2.8),
        java, size=14, color=OFF_WHITE, font=FONT_MONO,
    )
    _, tf = add_textbox(slide, Inches(0.7), Inches(5.2), Inches(11.5), Inches(1.8), "", size=16)
    add_bullets(
        tf,
        [
            "sa / empty password = local H2 default for teaching only.",
            "try-with-resources closes the connection.",
            "Verify: metadata query, H2 console, or INSERT + SELECT.",
            "“Table not found” → DDL not applied or wrong DB path.",
        ],
        size=16,
    )
    add_footer(slide)


def slide_step4_stores(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Step 4")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "Talk to tables from Java (preview)", size=34, bold=True, color=NAVY, font=FONT_TITLE,
    )
    rows = [
        ("TransactionRecordStore", "InMemory* today", "JDBC INSERT/SELECT → transaction_records"),
        ("IdempotencyRegistry", "InMemory*", "UNIQUE key + lookup"),
        ("AiDecisionStore", "InMemory*", "ai_decision_records (Capstone 5)"),
    ]
    table = slide.shapes.add_table(4, 3, Inches(0.7), Inches(1.85), Inches(11.9), Inches(2.8))
    tbl = table.table
    for j, h in enumerate(["Interface", "Today", "Capstone 4+"]):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            tbl.cell(i, j).text = val
            tbl.cell(i, j).text_frame.paragraphs[0].font.size = Pt(11)
    add_round_rect(slide, Inches(0.7), Inches(5.0), Inches(11.9), Inches(1.35), BLUE_PANEL)
    add_textbox(
        slide, Inches(1.0), Inches(5.25), Inches(11.3), Inches(0.9),
        "Round trip: Customer.name → customers.name → ResultSet.getString(\"name\")\n"
        "Steps 1–2 = class → column. Capstone 4 = full save/load + pipeline.",
        size=15, color=NAVY,
    )
    add_footer(slide)


def slide_step5_wipe(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Step 5")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "Reset the schema (local wipe)", size=34, bold=True, color=NAVY, font=FONT_TITLE,
    )
    _, tf = add_textbox(slide, Inches(0.7), Inches(1.85), Inches(6.5), Inches(4.5), "", size=17)
    add_bullets(
        tf,
        [
            "1. Stop app / Maven (release DB lock)",
            "2. Delete ./data/paynest.mv.db (and related H2 files)",
            "3. Restart → Step 3 re-runs DDL",
            "4. Re-run tests that need a clean store",
        ],
        size=17,
    )
    warn = add_round_rect(slide, Inches(7.5), Inches(2.0), Inches(5.1), Inches(3.5), NAVY)
    add_textbox(
        slide, Inches(7.8), Inches(2.3), Inches(4.5), Inches(0.5),
        "Learning vs production", size=16, bold=True, color=TEAL_LIGHT, font=FONT_TITLE,
    )
    _, tf2 = add_textbox(slide, Inches(7.8), Inches(2.9), Inches(4.5), Inches(2.3), "", size=14, color=OFF_WHITE)
    add_bullets(
        tf2,
        [
            "Wipe + recreate: fine on student laptop",
            "Later: schema-v2.sql with ALTER TABLE",
            "Never delete DB files while connection is open",
        ],
        size=14,
        color=OFF_WHITE,
    )
    add_footer(slide)


def slide_not_using(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Context")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "What we are not using (yet)", size=34, bold=True, color=NAVY, font=FONT_TITLE,
    )
    rows = [
        ("Spring ddl-auto", "Auto tables from entities", "Not required here"),
        ("Flyway / Liquibase", "Numbered migration runners", "Same SQL later if approved"),
        ("JPA / Hibernate", "Annotations map class ↔ table", "Plain JDBC in Capstone 4"),
        (".NET EF migrations add", "Generates migration from model", "You write SQL manually"),
    ]
    table = slide.shapes.add_table(5, 3, Inches(0.7), Inches(1.85), Inches(11.9), Inches(3.5))
    tbl = table.table
    for j, h in enumerate(["Tool", "Elsewhere", "PayNest"]):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            tbl.cell(i, j).text = val
            tbl.cell(i, j).text_frame.paragraphs[0].font.size = Pt(11)
    add_textbox(
        slide, Inches(0.7), Inches(5.6), Inches(11.5), Inches(0.8),
        "Hard part you are learning: clear DDL from a domain model. Runners are optional later.",
        size=16, color=SLATE,
    )
    add_footer(slide)


def slide_lab_wrap(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12), TEAL_LIGHT)
    add_textbox(
        slide, Inches(0.7), Inches(0.6), Inches(11), Inches(0.8),
        "Lab & wrap-up", size=34, bold=True, color=WHITE, font=FONT_TITLE,
    )
    _, tf = add_textbox(slide, Inches(0.9), Inches(1.6), Inches(11.5), Inches(4.2), "", size=18, color=OFF_WHITE)
    add_bullets(
        tf,
        [
            "List fields in TransactionRecord.java",
            "Draft CREATE TABLE + 3 constraints on paper",
            "Compare to schema-v1.sql",
            "Explain UNIQUE idempotency (duplicate webhook story)",
            "Write wipe runbook from memory",
            "Optional: open paynest-erd.drawio — match boxes to classes",
        ],
        size=18,
        color=OFF_WHITE,
        spacing=10,
    )
    add_round_rect(slide, Inches(0.7), Inches(5.9), Inches(11.9), Inches(0.85), TEAL)
    add_textbox(
        slide, Inches(1.0), Inches(6.05), Inches(11.3), Inches(0.55),
        "Step 1 domain → Step 2 migrations → Step 3 apply → stores & wipe",
        size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    slide.notes_slide.notes_text_frame.text = (
        "Point students to 01-code-first-migrations-walkthrough.md and Capstone 4 brief."
    )


def build():
    prs = new_presentation()
    slide_title(prs)
    slide_problem(prs)
    slide_session_spine(prs)
    slide_erd_visual_only(prs)
    slide_glossary_1(prs)
    slide_glossary_2(prs)
    slide_code_first(prs)
    slide_step1_intro(prs)
    slide_step1_nested(prs)
    slide_step1_mapping(prs)
    slide_step1_checkpoint(prs)
    slide_step2_intro(prs)
    slide_step2_sql_keywords(prs)
    slide_step2_create_table(prs)
    slide_step2_mistakes(prs)
    slide_step2_checkpoint(prs)
    slide_step3_apply(prs)
    slide_step3_jdbc(prs)
    slide_step4_stores(prs)
    slide_step5_wipe(prs)
    slide_not_using(prs)
    slide_lab_wrap(prs)
    prs.save(OUT_FILE)
    print(f"Wrote {OUT_FILE} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
