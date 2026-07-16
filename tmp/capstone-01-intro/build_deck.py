#!/usr/bin/env python3
"""Build Capstone 1 student introduction deck (Google Slides–compatible PPTX)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT_DIR = Path(__file__).resolve().parent
OUT_FILE = OUT_DIR / "capstone-01-intro.pptx"

# PayNest-inspired palette
NAVY = RGBColor(0x0F, 0x17, 0x2A)
TEAL = RGBColor(0x0D, 0x94, 0x88)
TEAL_LIGHT = RGBColor(0x14, 0xB8, 0xA6)
SLATE = RGBColor(0x47, 0x55, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF8, 0xFA, 0xFC)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)

FONT_TITLE = "Poppins"
FONT_BODY = "Lato"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line=False):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if not line:
        shape.line.fill.background()
    return shape


def add_round_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, size=18, bold=False,
                color=NAVY, align=PP_ALIGN.LEFT, font=FONT_BODY, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return box, tf


def add_bullets(tf, items, size=16, color=SLATE, spacing=6):
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT_BODY
        p.space_after = Pt(spacing)
        p.bullet = True


def add_header_bar(slide, kicker=None):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12), TEAL)
    if kicker:
        add_textbox(
            slide, Inches(0.7), Inches(0.35), Inches(4), Inches(0.35),
            kicker.upper(), size=11, bold=True, color=TEAL, font=FONT_TITLE
        )


def add_footer(slide, text="PayNest Capstone Programme"):
    add_textbox(
        slide, Inches(0.7), Inches(7.05), Inches(6), Inches(0.3),
        text, size=9, color=SLATE
    )


def slide_title(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, TEAL)
    add_round_rect(slide, Inches(9.2), Inches(1.2), Inches(3.6), Inches(5.1), TEAL_LIGHT)
    add_round_rect(slide, Inches(9.6), Inches(1.6), Inches(2.8), Inches(4.3), TEAL)

    add_textbox(
        slide, Inches(0.9), Inches(1.6), Inches(7.8), Inches(0.5),
        "CAPSTONE 1", size=14, bold=True, color=TEAL_LIGHT, font=FONT_TITLE
    )
    add_textbox(
        slide, Inches(0.9), Inches(2.1), Inches(7.8), Inches(1.6),
        "Merchant Order Desk\n& Catalogue Engine",
        size=40, bold=True, color=WHITE, font=FONT_TITLE
    )
    add_textbox(
        slide, Inches(0.9), Inches(4.0), Inches(7.5), Inches(1.0),
        "An introduction for students — build a trustworthy commerce kernel in plain Java.",
        size=20, color=OFF_WHITE, font=FONT_BODY
    )
    add_textbox(
        slide, Inches(0.9), Inches(6.2), Inches(7), Inches(0.4),
        "PayNest · Junior Backend Engineering",
        size=13, color=SLATE, font=FONT_BODY
    )
    slide.notes_slide.notes_text_frame.text = (
        "Welcome students. Set expectations: this is Capstone 1, focused on domain modelling "
        "and correct order totals — not payments yet."
    )


def slide_welcome(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Welcome")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "What is this capstone?", size=34, bold=True, color=NAVY, font=FONT_TITLE
    )
    add_textbox(
        slide, Inches(0.7), Inches(1.65), Inches(11), Inches(0.6),
        "Your first production-minded backend assignment in the PayNest programme.",
        size=18, color=SLATE
    )

    cards = [
        ("Goal", "Ship a minimal commerce kernel merchants can trust."),
        ("Focus", "Products, customers, orders, and correct totals."),
        ("Style", "Plain Java objects — no frameworks, no database."),
        ("Outcome", "A CLI demo that prints a clear order summary."),
    ]
    for i, (title, body) in enumerate(cards):
        col = i % 2
        row = i // 2
        left = Inches(0.7 + col * 6.1)
        top = Inches(2.6 + row * 2.2)
        card = add_round_rect(slide, left, top, Inches(5.7), Inches(1.85), WHITE)
        card.shadow.inherit = False
        add_rect(slide, left, top, Inches(0.08), Inches(1.85), TEAL)
        add_textbox(slide, left + Inches(0.35), top + Inches(0.25), Inches(5), Inches(0.4),
                    title, size=16, bold=True, color=TEAL, font=FONT_TITLE)
        add_textbox(slide, left + Inches(0.35), top + Inches(0.7), Inches(5), Inches(0.9),
                    body, size=15, color=SLATE)
    add_footer(slide)
    slide.notes_slide.notes_text_frame.text = (
        "Emphasise this is foundational — correctness and clear modelling matter more than polish."
    )


def slide_paynest_context(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Scenario")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "Meet PayNest", size=34, bold=True, color=NAVY, font=FONT_TITLE
    )

    panel = add_round_rect(slide, Inches(0.7), Inches(1.9), Inches(7.2), Inches(4.6), WHITE)
    _, tf = add_textbox(
        slide, Inches(1.0), Inches(2.15), Inches(6.6), Inches(4.2), "", size=16, color=SLATE
    )
    bullets = [
        "Early-stage South African fintech building lightweight commerce tools.",
        "Target customers: small merchants selling hardware and accessories online and at markets.",
        "They cannot afford Shopify-scale subscriptions.",
        "They still need consistent pricing, order totals, and customer-linked receipts.",
        "Payment integration comes later — first, prove the commerce kernel works.",
    ]
    add_bullets(tf, bullets, size=16)

    add_round_rect(slide, Inches(8.3), Inches(2.0), Inches(4.3), Inches(2.2), NAVY)
    add_textbox(
        slide, Inches(8.6), Inches(2.3), Inches(3.7), Inches(0.5),
        "Your brief", size=14, bold=True, color=TEAL_LIGHT, font=FONT_TITLE
    )
    add_textbox(
        slide, Inches(8.6), Inches(2.9), Inches(3.7), Inches(1.1),
        '"Ship something we can demo to merchants next week — no frameworks, just solid Java objects we can extend later."',
        size=14, color=OFF_WHITE
    )
    add_textbox(
        slide, Inches(8.6), Inches(4.5), Inches(3.7), Inches(0.5),
        "Currency: Rand (R)", size=13, bold=True, color=AMBER, font=FONT_TITLE
    )
    add_footer(slide)


def slide_business_problem(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Problem")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "The business problem", size=34, bold=True, color=NAVY, font=FONT_TITLE
    )
    add_textbox(
        slide, Inches(0.7), Inches(1.65), Inches(11), Inches(0.5),
        "Spreadsheets and WhatsApp messages break down as merchants grow.",
        size=18, color=SLATE
    )

    pains = [
        ("Line totals disagree", "with invoices and receipts"),
        ("Duplicate products", "get added by mistake"),
        ("No single code path", 'computes "what the customer owes"'),
    ]
    for i, (head, sub) in enumerate(pains):
        left = Inches(0.7 + i * 4.1)
        card = add_round_rect(slide, left, Inches(2.5), Inches(3.8), Inches(2.4), WHITE)
        add_textbox(slide, left + Inches(0.3), Inches(2.8), Inches(3.2), Inches(0.7),
                    head, size=17, bold=True, color=NAVY, font=FONT_TITLE)
        add_textbox(slide, left + Inches(0.3), Inches(3.5), Inches(3.2), Inches(0.8),
                    sub, size=14, color=SLATE)

    banner = add_round_rect(slide, Inches(0.7), Inches(5.3), Inches(11.9), Inches(1.2), TEAL)
    add_textbox(
        slide, Inches(1.0), Inches(5.55), Inches(11.3), Inches(0.8),
        "PayNest needs a minimal commerce kernel: products with prices, customers, orders with line items, and a trustworthy order total.",
        size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE
    )
    add_footer(slide)


def slide_your_role(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12), TEAL_LIGHT)
    add_textbox(
        slide, Inches(0.7), Inches(0.5), Inches(4), Inches(0.35),
        "YOUR ROLE", size=11, bold=True, color=TEAL_LIGHT, font=FONT_TITLE
    )
    add_textbox(
        slide, Inches(0.7), Inches(0.95), Inches(11), Inches(0.8),
        "Junior backend engineer on contract", size=34, bold=True, color=WHITE, font=FONT_TITLE
    )

    steps = [
        "Model the domain with small, cohesive Java classes.",
        "Implement correct line subtotals and grand totals.",
        "Print a human-readable order summary for demo reviewers.",
        "Write tests that would catch arithmetic or collection regressions.",
        "Explain how your design can grow without rewriting checkout.",
    ]
    for i, step in enumerate(steps):
        top = Inches(2.1 + i * 0.95)
        num = add_round_rect(slide, Inches(0.9), top, Inches(0.55), Inches(0.55), TEAL)
        add_textbox(
            slide, Inches(0.9), top + Inches(0.05), Inches(0.55), Inches(0.45),
            str(i + 1), size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
            font=FONT_TITLE, anchor=MSO_ANCHOR.MIDDLE
        )
        add_textbox(
            slide, Inches(1.65), top + Inches(0.05), Inches(10.5), Inches(0.55),
            step, size=17, color=OFF_WHITE
        )
    add_footer(slide, "Capstone 1 · Commerce Engine")


def slide_what_you_build(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Architecture")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "What you will build", size=34, bold=True, color=NAVY, font=FONT_TITLE
    )

    boxes = [
        ("Product", "id, name, price (R)"),
        ("Customer", "id, name, email"),
        ("OrderItem", "product × quantity → line total"),
        ("Order", "owns items, calculates grand total"),
        ("OrderService", "creates orders for customers"),
        ("PayNestApplication", "demo: products → order → summary"),
    ]
    positions = [
        (0.7, 2.0), (4.5, 2.0), (8.3, 2.0),
        (2.6, 4.0), (6.4, 4.0), (4.5, 5.8),
    ]
    for (label, sub), (x, y) in zip(boxes, positions):
        left = Inches(x)
        top = Inches(y)
        add_round_rect(slide, left, top, Inches(3.4), Inches(1.35), WHITE)
        add_rect(slide, left, top, Inches(3.4), Inches(0.06), TEAL)
        add_textbox(slide, left + Inches(0.2), top + Inches(0.2), Inches(3), Inches(0.45),
                    label, size=16, bold=True, color=NAVY, font=FONT_TITLE, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.2), top + Inches(0.65), Inches(3), Inches(0.5),
                    sub, size=12, color=SLATE, align=PP_ALIGN.CENTER)

  # arrows as simple connectors via thin rects
    add_rect(slide, Inches(4.1), Inches(2.65), Inches(0.35), Inches(0.04), TEAL)
    add_rect(slide, Inches(8.0), Inches(2.65), Inches(0.25), Inches(0.04), TEAL)
    add_rect(slide, Inches(3.5), Inches(3.35), Inches(0.04), Inches(0.6), TEAL)
    add_rect(slide, Inches(7.3), Inches(3.35), Inches(0.04), Inches(0.6), TEAL)
    add_rect(slide, Inches(5.5), Inches(5.35), Inches(0.04), Inches(0.4), TEAL)

    add_textbox(
        slide, Inches(0.7), Inches(6.55), Inches(11.5), Inches(0.4),
        "Read-only checkout for Capstone 1 — focus on modelling and correct totals.",
        size=13, color=SLATE, align=PP_ALIGN.CENTER
    )
    add_footer(slide)


def slide_domain_product_customer(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Requirements")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "Domain model: Product & Customer", size=32, bold=True, color=NAVY, font=FONT_TITLE
    )

    for i, (title, items) in enumerate([
        ("Product catalogue", [
            "Private fields: id, name, unit price (Rand).",
            'Constructor: new Product(1, "Laptop", 12000)',
            "Getters for name and price when building summaries.",
        ]),
        ("Customer identity", [
            "Private fields: id, name, contact (e.g. email).",
            "Constructor and getters for receipt header.",
            "Links a person to every order they place.",
        ]),
    ]):
        left = Inches(0.7 + i * 6.2)
        add_round_rect(slide, left, Inches(1.9), Inches(5.8), Inches(4.5), WHITE)
        add_textbox(slide, left + Inches(0.35), Inches(2.15), Inches(5.2), Inches(0.5),
                    title, size=20, bold=True, color=TEAL, font=FONT_TITLE)
        _, tf = add_textbox(slide, left + Inches(0.35), Inches(2.75), Inches(5.1), Inches(3.4), "")
        add_bullets(tf, items, size=15)
    add_footer(slide)


def slide_domain_order(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Requirements")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "Domain model: Order & OrderItem", size=32, bold=True, color=NAVY, font=FONT_TITLE
    )

    _, tf = add_textbox(slide, Inches(0.7), Inches(1.85), Inches(7.5), Inches(4.8), "")
    add_bullets(tf, [
        "OrderItem links one Product with a quantity; calculateTotal() = price × quantity.",
        "Order stores id, Customer, and List<OrderItem>.",
        "addItem(Product, quantity) — quantities must be positive integers.",
        "calculateTotal() loops items and sums line subtotals.",
        "printSummary() lists each line (name, qty, subtotal) plus grand total.",
        "Console output must let reviewers reconcile totals manually — no hidden magic.",
    ], size=16)

    add_round_rect(slide, Inches(8.6), Inches(2.0), Inches(4.0), Inches(4.2), NAVY)
    add_textbox(slide, Inches(8.9), Inches(2.25), Inches(3.4), Inches(0.4),
                "Example summary", size=13, bold=True, color=TEAL_LIGHT, font=FONT_TITLE)
    lines = [
        "Customer: Thabo M.",
        "─────────────────",
        "Laptop    × 1  R12,000",
        "Mouse     × 2  R   400",
        "─────────────────",
        "Grand total    R12,400",
    ]
    for i, line in enumerate(lines):
        add_textbox(
            slide, Inches(8.9), Inches(2.75 + i * 0.42), Inches(3.4), Inches(0.38),
            line, size=12, color=OFF_WHITE if i > 0 else TEAL_LIGHT,
            font="Courier New" if i > 0 else FONT_BODY
        )
    add_footer(slide)


def slide_implementation_order(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Getting started")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "Suggested implementation order", size=32, bold=True, color=NAVY, font=FONT_TITLE
    )
    add_textbox(
        slide, Inches(0.7), Inches(1.6), Inches(11), Inches(0.45),
        "Complete classes in this order before worrying about polish.",
        size=16, color=SLATE
    )

    steps = [
        "Create Product and Customer data classes.",
        "Create OrderItem; verify one line subtotal is correct.",
        "Create Order with an empty List<OrderItem> in the constructor.",
        "Add Order#addItem(...) and confirm the list grows.",
        "Add Order#calculateTotal() by summing each OrderItem total.",
        "Add Order#printSummary() only after calculations are correct.",
        "Wire the full flow in PayNestApplication.",
    ]
    for i, step in enumerate(steps):
        row_top = Inches(2.15 + i * 0.68)
        add_round_rect(slide, Inches(0.9), row_top, Inches(0.48), Inches(0.48), TEAL)
        add_textbox(
            slide, Inches(0.9), row_top + Inches(0.04), Inches(0.48), Inches(0.4),
            str(i + 1), size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
            font=FONT_TITLE, anchor=MSO_ANCHOR.MIDDLE
        )
        add_textbox(slide, Inches(1.55), row_top + Inches(0.06), Inches(10.5), Inches(0.45),
                    step, size=15, color=NAVY)
    add_footer(slide)


def slide_technical_constraints(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Constraints")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "Technical constraints & business rules", size=30, bold=True, color=NAVY, font=FONT_TITLE
    )

    constraints = [
        ("Stack", "Java 21 · Maven · JUnit 5"),
        ("Architecture", "Plain Java — no Spring, no app server, no database"),
        ("Packages", "com.paynestsystem — domain, service, app"),
        ("Collections", "JDK List for order lines; avoid parallel streams"),
        ("Line subtotal", "unitPrice × quantity (double arithmetic)"),
        ("Grand total", "Must equal sum of line subtotals shown in summary"),
        ("Quantities", "Must be > 0; reject or guard invalid adds consistently"),
        ("Encapsulation", "Don't expose mutable internals that corrupt totals"),
    ]
    for i, (label, value) in enumerate(constraints):
        col = i % 2
        row = i // 2
        left = Inches(0.7 + col * 6.1)
        top = Inches(1.85 + row * 1.25)
        add_round_rect(slide, left, top, Inches(5.7), Inches(1.0), WHITE)
        add_textbox(slide, left + Inches(0.25), top + Inches(0.12), Inches(1.6), Inches(0.35),
                    label, size=12, bold=True, color=TEAL, font=FONT_TITLE)
        add_textbox(slide, left + Inches(0.25), top + Inches(0.45), Inches(5.2), Inches(0.45),
                    value, size=14, color=NAVY)
    add_footer(slide)


def slide_deliverables(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Submission")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "Deliverables & demo workflow", size=32, bold=True, color=NAVY, font=FONT_TITLE
    )

    _, tf = add_textbox(slide, Inches(0.7), Inches(1.85), Inches(6.5), Inches(4.5), "")
    add_bullets(tf, [
        "Source code implementing the commerce flow in the repo layout.",
        "mvn test passes on your branch (add tests for totals and edge cases).",
        "Short setup note: how to run the demo and what output reviewers expect.",
        "Optional: simple diagram (OrderService → domain → summary).",
        "Submit per programme rules (zip, Git tag, or LMS upload).",
    ], size=16)

    add_round_rect(slide, Inches(7.6), Inches(1.9), Inches(5.0), Inches(4.4), NAVY)
    add_textbox(slide, Inches(7.9), Inches(2.15), Inches(4.4), Inches(0.4),
                "Demo checklist", size=14, bold=True, color=TEAL_LIGHT, font=FONT_TITLE)
    demo = [
        "≥ 2 sample Product objects",
        "1 sample Customer",
        "Order with ≥ 2 line items",
        "At least one qty > 1",
        "Printed summary with grand total",
    ]
    for i, item in enumerate(demo):
        add_textbox(
            slide, Inches(7.9), Inches(2.65 + i * 0.55), Inches(4.4), Inches(0.45),
            f"✓  {item}", size=14, color=OFF_WHITE
        )
    add_footer(slide)


def slide_rubric(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Assessment")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "How you will be assessed", size=32, bold=True, color=NAVY, font=FONT_TITLE
    )

    rows = [
        ("Architecture & domain modelling", "25%", "Clear separation; coherent OrderService; extensible design"),
        ("Correctness & business rules", "30%", "Matching subtotals; validated quantities; summary matches computation"),
        ("Testing & verification", "15%", "Tests cover totals and edge cases; mvn test green"),
        ("Code quality & maintainability", "20%", "Readable code, encapsulation, no dead paths in main"),
        ("Documentation & communication", "10%", "Setup instructions; comments on non-obvious rules"),
    ]

    # table header
    add_round_rect(slide, Inches(0.7), Inches(1.85), Inches(11.9), Inches(0.55), TEAL)
    for col, label, left, width in [
        (0, "Category", Inches(0.9), Inches(4.8)),
        (1, "Weight", Inches(5.8), Inches(1.0)),
        (2, "What reviewers look for", Inches(6.9), Inches(5.5)),
    ]:
        add_textbox(slide, left, Inches(1.95), width, Inches(0.4),
                    label, size=12, bold=True, color=WHITE, font=FONT_TITLE)

    for i, (cat, weight, detail) in enumerate(rows):
        top = Inches(2.5 + i * 0.85)
        bg = WHITE if i % 2 == 0 else RGBColor(0xF1, 0xF5, 0xF9)
        add_round_rect(slide, Inches(0.7), top, Inches(11.9), Inches(0.75), bg)
        add_textbox(slide, Inches(0.9), top + Inches(0.15), Inches(4.6), Inches(0.5),
                    cat, size=13, bold=True, color=NAVY)
        add_textbox(slide, Inches(5.8), top + Inches(0.15), Inches(0.9), Inches(0.5),
                    weight, size=13, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(6.9), top + Inches(0.12), Inches(5.5), Inches(0.55),
                    detail, size=12, color=SLATE)

    banner = add_round_rect(slide, Inches(0.7), Inches(6.0), Inches(11.9), Inches(0.75), NAVY)
    add_textbox(
        slide, Inches(1.0), Inches(6.15), Inches(11.3), Inches(0.5),
        "Pass expectation: no critical correctness failures in totals.",
        size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE
    )
    add_footer(slide)


def slide_closing(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12), TEAL_LIGHT)
    add_textbox(
        slide, Inches(0.7), Inches(1.5), Inches(11.5), Inches(1.2),
        "Ready to build?", size=40, bold=True, color=WHITE, font=FONT_TITLE, align=PP_ALIGN.CENTER
    )
    add_textbox(
        slide, Inches(1.5), Inches(2.9), Inches(10.3), Inches(1.0),
        "Start with Product and Customer. Get one line total right before you print a summary. Ask early if a requirement is unclear.",
        size=20, color=OFF_WHITE, align=PP_ALIGN.CENTER
    )

    add_round_rect(slide, Inches(3.5), Inches(4.5), Inches(6.3), Inches(1.5), TEAL)
    add_textbox(
        slide, Inches(3.7), Inches(4.75), Inches(5.9), Inches(1.0),
        "Read the full brief:\ndocs/assessments/capstone-01-commerce-engine.md",
        size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE
    )
    add_textbox(
        slide, Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.4),
        "Good luck — ship something merchants can trust.",
        size=14, color=SLATE, align=PP_ALIGN.CENTER
    )
    slide.notes_slide.notes_text_frame.text = (
        "Point students to the assessment doc and starter repo. Encourage incremental commits and tests."
    )


def build():
    prs = new_presentation()
    slide_title(prs)
    slide_welcome(prs)
    slide_paynest_context(prs)
    slide_business_problem(prs)
    slide_your_role(prs)
    slide_what_you_build(prs)
    slide_domain_product_customer(prs)
    slide_domain_order(prs)
    slide_implementation_order(prs)
    slide_technical_constraints(prs)
    slide_deliverables(prs)
    slide_rubric(prs)
    slide_closing(prs)
    prs.save(OUT_FILE)
    print(f"Saved {OUT_FILE} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
