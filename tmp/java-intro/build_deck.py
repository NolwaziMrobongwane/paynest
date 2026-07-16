#!/usr/bin/env python3
"""Build Introduction to Java slide deck (W3Schools reference)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT_DIR = Path(__file__).resolve().parent
OUT_FILE = OUT_DIR / "java-intro.pptx"

# W3Schools-inspired palette (green + dark slate)
GREEN = RGBColor(0x04, 0xAA, 0x6D)
GREEN_DARK = RGBColor(0x03, 0x7A, 0x4F)
GREEN_LIGHT = RGBColor(0x9F, 0xE2, 0xB3)
NAVY = RGBColor(0x28, 0x2A, 0x35)
SLATE = RGBColor(0x47, 0x55, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF8, 0xFA, 0xFC)
CODE_BG = RGBColor(0x1E, 0x29, 0x3B)
CODE_TEXT = RGBColor(0xE2, 0xE8, 0xF0)
ORANGE = RGBColor(0xFF, 0x98, 0x4A)

FONT_TITLE = "Poppins"
FONT_BODY = "Lato"
FONT_CODE = "Courier New"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_round_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, size=18, bold=False,
                color=NAVY, align=PP_ALIGN.LEFT, font=FONT_BODY,
                anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return box, tf


def add_bullets(tf, items, size=16, color=SLATE, spacing=6, level=0):
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = level
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT_BODY
        p.space_after = Pt(spacing)
        if level == 0:
            p.bullet = True


def add_header_bar(slide, kicker=None):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12), GREEN)
    if kicker:
        add_textbox(
            slide, Inches(0.7), Inches(0.35), Inches(5), Inches(0.35),
            kicker.upper(), size=11, bold=True, color=GREEN, font=FONT_TITLE
        )


def add_footer(slide, text="Reference: W3Schools Java Tutorial"):
    add_textbox(
        slide, Inches(0.7), Inches(7.05), Inches(8), Inches(0.3),
        text, size=9, color=SLATE
    )


def add_code_block(slide, left, top, width, height, lines):
    add_round_rect(slide, left, top, width, height, CODE_BG)
    y = top + Inches(0.25)
    for line in lines:
        add_textbox(
            slide, left + Inches(0.3), y, width - Inches(0.5), Inches(0.38),
            line, size=13, color=CODE_TEXT, font=FONT_CODE
        )
        y += Inches(0.36)


def slide_title(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, GREEN)
    add_round_rect(slide, Inches(9.0), Inches(1.0), Inches(3.8), Inches(5.5), GREEN)
    add_round_rect(slide, Inches(9.5), Inches(1.5), Inches(2.8), Inches(4.5), GREEN_DARK)

    add_textbox(
        slide, Inches(0.9), Inches(1.5), Inches(7.8), Inches(0.5),
        "PROGRAMMING FUNDAMENTALS", size=13, bold=True, color=GREEN_LIGHT, font=FONT_TITLE
    )
    add_textbox(
        slide, Inches(0.9), Inches(2.0), Inches(7.8), Inches(1.4),
        "Introduction to Java",
        size=44, bold=True, color=WHITE, font=FONT_TITLE
    )
    add_textbox(
        slide, Inches(0.9), Inches(3.6), Inches(7.5), Inches(1.0),
        "A beginner-friendly overview based on the W3Schools Java tutorial.",
        size=20, color=OFF_WHITE
    )
    add_textbox(
        slide, Inches(0.9), Inches(6.2), Inches(7), Inches(0.4),
        "w3schools.com/java/java_intro.asp",
        size=12, color=SLATE
    )
    slide.notes_slide.notes_text_frame.text = (
        "Welcome. This deck introduces Java using W3Schools as the reference point. "
        "Encourage students to follow the tutorial in order."
    )


def slide_what_is_java(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Overview")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "What is Java?", size=36, bold=True, color=NAVY, font=FONT_TITLE
    )

    add_round_rect(slide, Inches(0.7), Inches(1.9), Inches(7.4), Inches(2.2), WHITE)
    add_textbox(
        slide, Inches(1.0), Inches(2.15), Inches(6.8), Inches(1.7),
        "Java is a popular and powerful programming language, created in 1995. "
        "It is owned by Oracle, and more than 3 billion devices run Java worldwide.",
        size=18, color=SLATE
    )

    stats = [
        ("1995", "Language created"),
        ("3B+", "Devices run Java"),
        ("Oracle", "Current owner"),
    ]
    for i, (num, label) in enumerate(stats):
        left = Inches(8.4 + (i % 2) * 2.3)
        top = Inches(1.9 + (i // 2) * 2.3) if i < 2 else Inches(4.2)
        if i == 2:
            left = Inches(9.55)
        add_round_rect(slide, left, top, Inches(2.1), Inches(1.8), NAVY)
        add_textbox(
            slide, left + Inches(0.15), top + Inches(0.35), Inches(1.8), Inches(0.6),
            num, size=26, bold=True, color=GREEN_LIGHT, align=PP_ALIGN.CENTER, font=FONT_TITLE
        )
        add_textbox(
            slide, left + Inches(0.15), top + Inches(1.05), Inches(1.8), Inches(0.5),
            label, size=12, color=OFF_WHITE, align=PP_ALIGN.CENTER
        )

    add_textbox(
        slide, Inches(0.7), Inches(4.5), Inches(11), Inches(0.4),
        "Java is used across many domains — not just one type of application.",
        size=16, bold=True, color=NAVY
    )
    add_footer(slide)


def slide_where_used(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Applications")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "Where is Java used?", size=34, bold=True, color=NAVY, font=FONT_TITLE
    )

    uses = [
        ("Mobile", "Android apps"),
        ("Desktop", "Cross-platform apps"),
        ("Web", "Web applications"),
        ("Servers", "Application servers"),
        ("Games", "Game development"),
        ("Data", "Database connections"),
    ]
    for i, (title, sub) in enumerate(uses):
        col = i % 3
        row = i // 3
        left = Inches(0.7 + col * 4.1)
        top = Inches(2.0 + row * 2.3)
        add_round_rect(slide, left, top, Inches(3.8), Inches(1.9), WHITE)
        add_rect(slide, left, top, Inches(3.8), Inches(0.07), GREEN)
        add_textbox(
            slide, left + Inches(0.3), top + Inches(0.35), Inches(3.2), Inches(0.5),
            title, size=20, bold=True, color=GREEN_DARK, font=FONT_TITLE
        )
        add_textbox(
            slide, left + Inches(0.3), top + Inches(0.95), Inches(3.2), Inches(0.6),
            sub, size=15, color=SLATE
        )
    add_footer(slide)


def slide_why_java(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Benefits")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "Why use Java?", size=34, bold=True, color=NAVY, font=FONT_TITLE
    )

    reasons = [
        "Works on different platforms (Windows, Mac, Linux, Raspberry Pi, …)",
        "One of the most popular programming languages in the world",
        "Strong demand in the current job market",
        "Easy to learn and simple to use for beginners",
        "Open-source and free",
        "Secure, fast, and powerful",
        "Huge community support (tens of millions of developers)",
        "Object-oriented — clear program structure and reusable code",
    ]
    for i, reason in enumerate(reasons):
        col = i % 2
        row = i // 2
        left = Inches(0.7 + col * 6.1)
        top = Inches(1.85 + row * 1.2)
        add_round_rect(slide, left, top, Inches(5.7), Inches(0.95), WHITE)
        add_textbox(
            slide, left + Inches(0.3), top + Inches(0.22), Inches(5.1), Inches(0.55),
            reason, size=14, color=NAVY
        )

    banner = add_round_rect(slide, Inches(0.7), Inches(6.0), Inches(11.9), Inches(0.7), GREEN)
    add_textbox(
        slide, Inches(1.0), Inches(6.12), Inches(11.3), Inches(0.45),
        "Java is close to C++ and C# — easy to switch between these languages.",
        size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE
    )
    add_footer(slide)


def slide_platform_independence(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Concept")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "Write once, run anywhere", size=34, bold=True, color=NAVY, font=FONT_TITLE
    )
    add_textbox(
        slide, Inches(0.7), Inches(1.65), Inches(11), Inches(0.5),
        "Java source code is compiled to bytecode, then runs on the Java Virtual Machine (JVM).",
        size=17, color=SLATE
    )

    steps = [
        ("1. Write", ".java source file"),
        ("2. Compile", "javac → .class bytecode"),
        ("3. Run", "JVM on any platform"),
    ]
    for i, (title, sub) in enumerate(steps):
        left = Inches(0.9 + i * 4.0)
        add_round_rect(slide, left, Inches(2.5), Inches(3.5), Inches(2.0), WHITE)
        add_textbox(
            slide, left + Inches(0.25), Inches(2.8), Inches(3), Inches(0.5),
            title, size=20, bold=True, color=GREEN, font=FONT_TITLE, align=PP_ALIGN.CENTER
        )
        add_textbox(
            slide, left + Inches(0.25), Inches(3.45), Inches(3), Inches(0.7),
            sub, size=14, color=SLATE, align=PP_ALIGN.CENTER
        )
        if i < 2:
            add_rect(slide, left + Inches(3.55), Inches(3.35), Inches(0.35), Inches(0.05), GREEN)

    platforms = ["Windows", "macOS", "Linux", "Android", "Cloud servers"]
    add_round_rect(slide, Inches(0.7), Inches(5.0), Inches(11.9), Inches(1.5), NAVY)
    add_textbox(
        slide, Inches(1.0), Inches(5.25), Inches(11.3), Inches(0.4),
        "Same bytecode runs on many platforms", size=14, bold=True,
        color=GREEN_LIGHT, align=PP_ALIGN.CENTER, font=FONT_TITLE
    )
    for i, plat in enumerate(platforms):
        add_round_rect(slide, Inches(1.0 + i * 2.3), Inches(5.75), Inches(2.0), Inches(0.5), GREEN_DARK)
        add_textbox(
            slide, Inches(1.0 + i * 2.3), Inches(5.8), Inches(2.0), Inches(0.4),
            plat, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE
        )
    add_footer(slide)


def slide_oop_intro(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Paradigm")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "Java is object-oriented", size=34, bold=True, color=NAVY, font=FONT_TITLE
    )

    add_round_rect(slide, Inches(0.7), Inches(1.9), Inches(5.5), Inches(4.5), WHITE)
    _, tf = add_textbox(slide, Inches(1.0), Inches(2.15), Inches(5.0), Inches(4.0), "")
    add_bullets(tf, [
        "Programs are built from classes and objects.",
        "A class is a blueprint; an object is a real instance.",
        "OOP gives programs a clear structure.",
        "Code can be reused, lowering development costs.",
        "You'll learn classes, objects, methods, and encapsulation in depth.",
    ], size=16)

    add_round_rect(slide, Inches(6.6), Inches(2.0), Inches(6.0), Inches(1.6), CODE_BG)
    add_textbox(
        slide, Inches(6.9), Inches(2.15), Inches(5.4), Inches(0.35),
        "Class (blueprint)", size=12, bold=True, color=GREEN_LIGHT, font=FONT_TITLE
    )
    add_textbox(
        slide, Inches(6.9), Inches(2.55), Inches(5.4), Inches(0.9),
        'class Car { String brand; void honk() { … } }',
        size=13, color=CODE_TEXT, font=FONT_CODE
    )

    add_round_rect(slide, Inches(6.6), Inches(4.0), Inches(6.0), Inches(1.6), GREEN_DARK)
    add_textbox(
        slide, Inches(6.9), Inches(4.15), Inches(5.4), Inches(0.35),
        "Object (instance)", size=12, bold=True, color=WHITE, font=FONT_TITLE
    )
    add_textbox(
        slide, Inches(6.9), Inches(4.55), Inches(5.4), Inches(0.9),
        'Car myCar = new Car();\nmyCar.brand = "Toyota";',
        size=13, color=OFF_WHITE, font=FONT_CODE
    )
    add_footer(slide)


def slide_program_structure(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Syntax")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "Anatomy of a Java program", size=32, bold=True, color=NAVY, font=FONT_TITLE
    )

    add_code_block(slide, Inches(0.7), Inches(1.85), Inches(6.2), Inches(4.5), [
        "public class Main {",
        "  public static void main(String[] args) {",
        '    String name = "John";',
        '    System.out.println("Hello " + name);',
        "  }",
        "}",
    ])

    labels = [
        ("public class Main", "Every Java file needs a class. Filename must match."),
        ("main method", "Entry point — program starts here."),
        ("String name = …", "Variable declaration and assignment."),
        ("System.out.println", "Prints output to the console."),
    ]
    for i, (part, desc) in enumerate(labels):
        top = Inches(2.0 + i * 1.05)
        add_round_rect(slide, Inches(7.3), top, Inches(5.3), Inches(0.85), WHITE)
        add_textbox(
            slide, Inches(7.55), top + Inches(0.1), Inches(4.8), Inches(0.3),
            part, size=13, bold=True, color=GREEN_DARK, font=FONT_CODE
        )
        add_textbox(
            slide, Inches(7.55), top + Inches(0.42), Inches(4.8), Inches(0.35),
            desc, size=12, color=SLATE
        )
    add_footer(slide)


def slide_hello_world(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12), GREEN)
    add_textbox(
        slide, Inches(0.7), Inches(0.5), Inches(4), Inches(0.35),
        "EXAMPLE", size=11, bold=True, color=GREEN_LIGHT, font=FONT_TITLE
    )
    add_textbox(
        slide, Inches(0.7), Inches(0.95), Inches(11), Inches(0.8),
        "Your first Java program", size=34, bold=True, color=WHITE, font=FONT_TITLE
    )

    add_code_block(slide, Inches(0.9), Inches(2.0), Inches(7.5), Inches(3.2), [
        "public class Main {",
        "  public static void main(String[] args) {",
        '    String name = "John";',
        '    System.out.println("Hello " + name);',
        "  }",
        "}",
    ])

    add_round_rect(slide, Inches(8.8), Inches(2.2), Inches(3.8), Inches(1.4), GREEN)
    add_textbox(
        slide, Inches(9.0), Inches(2.45), Inches(3.4), Inches(0.4),
        "Output", size=13, bold=True, color=WHITE, font=FONT_TITLE, align=PP_ALIGN.CENTER
    )
    add_textbox(
        slide, Inches(9.0), Inches(2.95), Inches(3.4), Inches(0.45),
        "Hello John", size=20, bold=True, color=NAVY, font=FONT_CODE, align=PP_ALIGN.CENTER
    )

    add_textbox(
        slide, Inches(0.9), Inches(5.6), Inches(11.5), Inches(0.8),
        "This is the classic W3Schools example: declare a name, concatenate strings, and print to the console.",
        size=16, color=OFF_WHITE
    )
    add_footer(slide, "Source: w3schools.com/java/java_intro.asp")


def slide_key_concepts(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Roadmap")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "Core topics you will learn", size=32, bold=True, color=NAVY, font=FONT_TITLE
    )
    add_textbox(
        slide, Inches(0.7), Inches(1.6), Inches(11), Inches(0.45),
        "Follow the tutorial in order — each chapter builds on the previous one.",
        size=16, color=SLATE
    )

    topics = [
        ("Basics", "Syntax, output, comments, variables, data types"),
        ("Logic", "Operators, if/else, switch, loops"),
        ("Data", "Strings, arrays, methods"),
        ("OOP", "Classes, objects, constructors, encapsulation"),
        ("Advanced", "Inheritance, interfaces, exceptions, collections"),
        ("Practice", "Exercises, quizzes, and code challenges"),
    ]
    for i, (title, sub) in enumerate(topics):
        col = i % 3
        row = i // 3
        left = Inches(0.7 + col * 4.1)
        top = Inches(2.3 + row * 2.2)
        add_round_rect(slide, left, top, Inches(3.8), Inches(1.85), WHITE)
        num = add_round_rect(slide, left + Inches(0.25), top + Inches(0.25), Inches(0.45), Inches(0.45), GREEN)
        add_textbox(
            slide, left + Inches(0.25), top + Inches(0.28), Inches(0.45), Inches(0.4),
            str(i + 1), size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
            font=FONT_TITLE, anchor=MSO_ANCHOR.MIDDLE
        )
        add_textbox(
            slide, left + Inches(0.85), top + Inches(0.25), Inches(2.7), Inches(0.4),
            title, size=17, bold=True, color=NAVY, font=FONT_TITLE
        )
        add_textbox(
            slide, left + Inches(0.25), top + Inches(0.8), Inches(3.3), Inches(0.85),
            sub, size=13, color=SLATE
        )
    add_footer(slide)


def slide_variables_preview(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Preview")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "Variables & data types (preview)", size=30, bold=True, color=NAVY, font=FONT_TITLE
    )

    add_code_block(slide, Inches(0.7), Inches(1.85), Inches(5.5), Inches(3.8), [
        "int age = 25;",
        "double price = 19.99;",
        "boolean active = true;",
        'char grade = \'A\';',
        'String city = "Cape Town";',
    ])

    types = [
        ("int", "Whole numbers"),
        ("double", "Decimal numbers"),
        ("boolean", "true or false"),
        ("char", "Single character"),
        ("String", "Text (reference type)"),
    ]
    for i, (typ, desc) in enumerate(types):
        top = Inches(1.95 + i * 0.72)
        add_round_rect(slide, Inches(6.6), top, Inches(5.9), Inches(0.58), WHITE)
        add_textbox(
            slide, Inches(6.85), top + Inches(0.1), Inches(1.2), Inches(0.35),
            typ, size=13, bold=True, color=GREEN_DARK, font=FONT_CODE
        )
        add_textbox(
            slide, Inches(8.1), top + Inches(0.1), Inches(4.2), Inches(0.35),
            desc, size=13, color=SLATE
        )

    add_textbox(
        slide, Inches(0.7), Inches(6.0), Inches(11.5), Inches(0.5),
        "W3Schools covers casting, operators, and real-life examples in the Variables and Data Types sections.",
        size=14, color=SLATE
    )
    add_footer(slide)


def slide_control_flow_preview(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Preview")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "Control flow (preview)", size=32, bold=True, color=NAVY, font=FONT_TITLE
    )

    for i, (title, code) in enumerate([
        ("if / else", 'if (score >= 50) {\n  System.out.println("Pass");\n} else {\n  System.out.println("Fail");\n}'),
        ("for loop", 'for (int i = 0; i < 5; i++) {\n  System.out.println(i);\n}'),
        ("while loop", 'while (count > 0) {\n  count--;\n}'),
    ]):
        left = Inches(0.7 + i * 4.1)
        add_textbox(
            slide, left, Inches(1.75), Inches(3.8), Inches(0.4),
            title, size=18, bold=True, color=GREEN_DARK, font=FONT_TITLE
        )
        lines = code.split("\n")
        add_code_block(slide, left, Inches(2.2), Inches(3.8), Inches(2.8), lines)

    add_round_rect(slide, Inches(0.7), Inches(5.4), Inches(11.9), Inches(1.1), GREEN_DARK)
    add_textbox(
        slide, Inches(1.0), Inches(5.65), Inches(11.3), Inches(0.7),
        "Programs need decisions and repetition. W3Schools teaches if/else, switch, while, for, and for-each loops with real-life examples.",
        size=15, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE
    )
    add_footer(slide)


def slide_learning_path(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, OFF_WHITE)
    add_header_bar(slide, "Resources")
    add_textbox(
        slide, Inches(0.7), Inches(0.85), Inches(11), Inches(0.8),
        "How to get started", size=34, bold=True, color=NAVY, font=FONT_TITLE
    )

    _, tf = add_textbox(slide, Inches(0.7), Inches(1.85), Inches(6.8), Inches(4.5), "")
    add_bullets(tf, [
        "You don't need prior programming experience — just curiosity and practice.",
        "Follow the W3Schools Java tutorial in order.",
        "Each chapter builds on the previous one.",
        "Use the online Try It editor to run examples in your browser.",
        "Complete exercises and code challenges to reinforce learning.",
        "By the end, you'll write basic Java programs and apply skills to real examples.",
    ], size=16)

    add_round_rect(slide, Inches(7.8), Inches(2.0), Inches(4.8), Inches(3.8), NAVY)
    add_textbox(
        slide, Inches(8.1), Inches(2.25), Inches(4.2), Inches(0.4),
        "Start here", size=14, bold=True, color=GREEN_LIGHT, font=FONT_TITLE
    )
    links = [
        "Java Introduction",
        "Java Get Started",
        "Java Syntax",
        "Java Variables",
        "Java OOP",
    ]
    for i, link in enumerate(links):
        add_textbox(
            slide, Inches(8.1), Inches(2.75 + i * 0.55), Inches(4.2), Inches(0.45),
            f"→  {link}", size=14, color=OFF_WHITE
        )
    add_textbox(
        slide, Inches(8.1), Inches(5.35), Inches(4.2), Inches(0.4),
        "w3schools.com/java/",
        size=11, color=GREEN_LIGHT
    )
    add_footer(slide)


def slide_closing(prs):
    slide = blank_slide(prs)
    set_slide_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.12), GREEN)
    add_textbox(
        slide, Inches(0.7), Inches(1.4), Inches(11.9), Inches(1.0),
        "Ready to code?", size=40, bold=True, color=WHITE, font=FONT_TITLE, align=PP_ALIGN.CENTER
    )
    add_textbox(
        slide, Inches(1.5), Inches(2.6), Inches(10.3), Inches(1.0),
        "Open the W3Schools Java tutorial, run the Hello World example, and experiment by changing the name variable.",
        size=19, color=OFF_WHITE, align=PP_ALIGN.CENTER
    )

    add_round_rect(slide, Inches(3.8), Inches(4.2), Inches(5.7), Inches(1.4), GREEN)
    add_textbox(
        slide, Inches(4.0), Inches(4.45), Inches(5.3), Inches(0.9),
        "w3schools.com/java/java_intro.asp",
        size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE
    )
    add_textbox(
        slide, Inches(0.7), Inches(6.3), Inches(11.9), Inches(0.4),
        "Practice daily. Small steps compound into real programming skill.",
        size=14, color=SLATE, align=PP_ALIGN.CENTER
    )
    slide.notes_slide.notes_text_frame.text = (
        "Direct students to W3Schools Java Introduction and Get Started pages. "
        "Assign Hello World as first hands-on exercise."
    )


def build():
    prs = new_presentation()
    slide_title(prs)
    slide_what_is_java(prs)
    slide_where_used(prs)
    slide_why_java(prs)
    slide_platform_independence(prs)
    slide_oop_intro(prs)
    slide_program_structure(prs)
    slide_hello_world(prs)
    slide_key_concepts(prs)
    slide_variables_preview(prs)
    slide_control_flow_preview(prs)
    slide_learning_path(prs)
    slide_closing(prs)
    prs.save(OUT_FILE)
    print(f"Saved {OUT_FILE} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
